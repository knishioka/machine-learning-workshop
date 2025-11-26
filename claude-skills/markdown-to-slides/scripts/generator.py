"""
Presentation Generator - スライドデータからPPTXを生成

python-pptxライブラリを使用してPowerPointファイルを生成
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from parser import PresentationData, SlideContent


# テーマ設定
THEMES = {
    "default": {
        "title_color": RGBColor(0x00, 0x00, 0x00),
        "text_color": RGBColor(0x33, 0x33, 0x33),
        "accent_color": RGBColor(0x00, 0x66, 0xCC),
    },
    "dark": {
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0xCC, 0xCC, 0xCC),
        "accent_color": RGBColor(0x00, 0xAA, 0xFF),
    },
    "corporate": {
        "title_color": RGBColor(0x1A, 0x1A, 0x2E),
        "text_color": RGBColor(0x16, 0x21, 0x3E),
        "accent_color": RGBColor(0xE9, 0x4B, 0x3C),
    },
}

# アスペクト比の寸法（EMU単位）
ASPECT_RATIOS = {
    "16:9": (Inches(13.333), Inches(7.5)),
    "4:3": (Inches(10), Inches(7.5)),
}


class PresentationGenerator:
    """構造化データからPowerPointプレゼンテーションを生成"""

    def __init__(self, theme: str = "default", aspect: str = "16:9"):
        """テーマとアスペクト比を指定して初期化"""
        self.theme = THEMES.get(theme, THEMES["default"])
        self.width, self.height = ASPECT_RATIOS.get(aspect, ASPECT_RATIOS["16:9"])

    def create_presentation(
        self,
        data: PresentationData,
        output_path: Path
    ) -> None:
        """パースされたデータからPowerPointを生成"""
        prs = Presentation()
        prs.slide_width = self.width
        prs.slide_height = self.height

        for slide_data in data.slides:
            if slide_data.is_title_slide:
                self._add_title_slide(prs, slide_data)
            else:
                self._add_content_slide(prs, slide_data)

        # 空のプレゼンテーションの場合
        if not data.slides:
            self._add_empty_slide(prs, data.title)

        prs.save(str(output_path))

    def _add_title_slide(self, prs: Presentation, slide: SlideContent) -> None:
        """タイトルスライドを追加"""
        layout = prs.slide_layouts[6]  # 空白レイアウト
        pptx_slide = prs.slides.add_slide(layout)

        # 中央にタイトルを配置
        left = Inches(1)
        top = Inches(3)
        width = self.width - Inches(2)
        height = Inches(1.5)

        title_box = pptx_slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.paragraphs[0].text = slide.title
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs: Presentation, slide: SlideContent) -> None:
        """コンテンツスライドを追加"""
        layout = prs.slide_layouts[6]  # 空白レイアウト
        pptx_slide = prs.slides.add_slide(layout)

        # タイトルを追加
        if slide.title:
            title_box = pptx_slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5),
                self.width - Inches(1), Inches(1)
            )
            tf = title_box.text_frame
            tf.paragraphs[0].text = slide.title
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].font.bold = True

        # 箇条書きを追加
        if slide.bullets:
            content_box = pptx_slide.shapes.add_textbox(
                Inches(0.75), Inches(1.75),
                self.width - Inches(1.5), self.height - Inches(2.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(slide.bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                # インデントを処理
                indent_level = len(bullet) - len(bullet.lstrip())
                p.level = indent_level // 2
                p.text = f"• {bullet.strip()}"
                p.font.size = Pt(24 - (p.level * 2))

        # コードブロックを追加
        if slide.code_blocks:
            self._add_code_block(pptx_slide, slide.code_blocks[0])

    def _add_code_block(self, pptx_slide, code: str) -> None:
        """コードブロックをスライドに追加"""
        code_box = pptx_slide.shapes.add_textbox(
            Inches(0.75), Inches(4),
            self.width - Inches(1.5), Inches(2.5)
        )
        tf = code_box.text_frame
        tf.paragraphs[0].text = code
        tf.paragraphs[0].font.name = "Courier New"
        tf.paragraphs[0].font.size = Pt(14)

    def _add_empty_slide(self, prs: Presentation, title: str) -> None:
        """空のプレゼンテーション用プレースホルダースライド"""
        layout = prs.slide_layouts[6]
        pptx_slide = prs.slides.add_slide(layout)

        title_box = pptx_slide.shapes.add_textbox(
            Inches(1), Inches(3),
            self.width - Inches(2), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title or "Empty Presentation"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

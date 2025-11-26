"""
markdown-to-slides 最小限テストスイート

5テストでコア機能をカバー:
- Parser: 2テスト（タイトル抽出、スライド分割）
- Generator: 2テスト（ファイル生成、スライド数）
- Integration: 1テスト（エンドツーエンド）
"""
import pytest
import sys
import tempfile
from pathlib import Path

# scriptsディレクトリをパスに追加
sys.path.insert(0, str(Path(__file__).parent.parent / "scripts"))

from parser import MarkdownParser, PresentationData, SlideContent
from generator import PresentationGenerator


# =============================================================================
# Fixtures
# =============================================================================

@pytest.fixture
def sample_markdown() -> str:
    """テスト用の標準Markdownコンテンツ"""
    return """# プレゼンテーションタイトル

## はじめに
- ポイント1
- ポイント2

## まとめ
- 結論A
- 結論B
"""


@pytest.fixture
def temp_dir():
    """一時ディレクトリを提供"""
    with tempfile.TemporaryDirectory() as dir_path:
        yield Path(dir_path)


# =============================================================================
# Parser Tests (2)
# =============================================================================

class TestParser:
    """Markdownパーサーのテスト"""

    def test_extracts_h1_as_title(self):
        """H1見出しがプレゼンテーションタイトルになること"""
        parser = MarkdownParser()
        content = "# 私のプレゼン\n\n内容"
        result = parser.parse(content)

        assert result.title == "私のプレゼン"
        assert result.slides[0].is_title_slide is True

    def test_h2_creates_content_slides(self):
        """H2見出しごとにコンテンツスライドが作成されること"""
        parser = MarkdownParser()
        content = "# タイトル\n\n## スライド1\n内容1\n\n## スライド2\n内容2"
        result = parser.parse(content)

        # タイトルスライド + コンテンツスライド2枚
        assert len(result.slides) == 3
        assert result.slides[1].title == "スライド1"
        assert result.slides[2].title == "スライド2"


# =============================================================================
# Generator Tests (2)
# =============================================================================

class TestGenerator:
    """PPTX生成のテスト"""

    def test_creates_pptx_file(self, temp_dir):
        """有効なPPTXファイルが作成されること"""
        generator = PresentationGenerator()
        data = PresentationData(
            title="テスト",
            slides=[SlideContent(title="テストタイトル", is_title_slide=True)]
        )
        output_path = temp_dir / "output.pptx"

        generator.create_presentation(data, output_path)

        assert output_path.exists()
        assert output_path.stat().st_size > 0

    def test_creates_correct_slide_count(self, temp_dir):
        """正しいスライド数が生成されること"""
        from pptx import Presentation as PptxPresentation

        generator = PresentationGenerator()
        data = PresentationData(
            title="テスト",
            slides=[
                SlideContent(title="タイトル", is_title_slide=True),
                SlideContent(title="スライド1", bullets=["項目1"]),
                SlideContent(title="スライド2", bullets=["項目2"]),
            ]
        )
        output_path = temp_dir / "output.pptx"

        generator.create_presentation(data, output_path)

        # 生成されたPPTXを開いてスライド数を確認
        prs = PptxPresentation(str(output_path))
        assert len(prs.slides) == 3


# =============================================================================
# Integration Test (1)
# =============================================================================

class TestIntegration:
    """エンドツーエンドテスト"""

    def test_full_conversion(self, sample_markdown, temp_dir):
        """Markdownから完全なPPTXへの変換"""
        from pptx import Presentation as PptxPresentation

        # Markdownファイルを作成
        md_path = temp_dir / "test.md"
        md_path.write_text(sample_markdown, encoding="utf-8")

        # パース
        parser = MarkdownParser()
        data = parser.parse_file(md_path)

        # 生成
        output_path = temp_dir / "output.pptx"
        generator = PresentationGenerator()
        generator.create_presentation(data, output_path)

        # 検証
        assert output_path.exists()
        prs = PptxPresentation(str(output_path))

        # タイトルスライド + 2つのコンテンツスライド
        assert len(prs.slides) == 3
        assert data.title == "プレゼンテーションタイトル"


# =============================================================================
# Run Tests
# =============================================================================

if __name__ == "__main__":
    pytest.main([__file__, "-v"])
